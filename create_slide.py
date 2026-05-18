from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(blank_slide_layout)

# Set background to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# ===== WORD CLOUD SECTION (Top 60-70%) =====
# Add word cloud text box (occupies top portion)
word_cloud_left = Inches(0.5)
word_cloud_top = Inches(0.5)
word_cloud_width = Inches(9)
word_cloud_height = Inches(4.2)

word_cloud_box = slide.shapes.add_textbox(word_cloud_left, word_cloud_top, word_cloud_width, word_cloud_height)
word_cloud_frame = word_cloud_box.text_frame
word_cloud_frame.word_wrap = True
word_cloud_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Add editable word cloud text
word_cloud_text = word_cloud_frame.paragraphs[0]
word_cloud_text.text = "Digital Transformation  Cloud Migration  Data Analytics  Customer Experience  Automation  Scalability  Integration  Security  AI Innovation  Efficiency  Agility"
word_cloud_text.alignment = PP_ALIGN.CENTER
word_cloud_text.font.size = Pt(28)
word_cloud_text.font.bold = True
word_cloud_text.font.color.rgb = RGBColor(40, 60, 80)  # Dark slate blue
word_cloud_text.font.name = 'Segoe UI'

# Add subtle border to word cloud section
word_cloud_box.line.color.rgb = RGBColor(220, 220, 220)
word_cloud_box.line.width = Pt(1)
word_cloud_box.fill.background()

# ===== TWO-COLUMN SECTION (Bottom ~30%) =====
# Left Column: Business Needs
left_col_left = Inches(0.5)
left_col_top = Inches(4.7)
left_col_width = Inches(4.4)
left_col_height = Inches(2.5)

# Left column heading
left_heading_box = slide.shapes.add_textbox(left_col_left, left_col_top, left_col_width, Inches(0.4))
left_heading_frame = left_heading_box.text_frame
left_heading_para = left_heading_frame.paragraphs[0]
left_heading_para.text = "Business Needs"
left_heading_para.font.size = Pt(18)
left_heading_para.font.bold = True
left_heading_para.font.color.rgb = RGBColor(40, 60, 80)
left_heading_para.font.name = 'Segoe UI'

# Left column content area
left_content_box = slide.shapes.add_textbox(left_col_left, left_col_top + Inches(0.5), left_col_width, Inches(1.8))
left_content_frame = left_content_box.text_frame
left_content_frame.word_wrap = True
left_content_frame.vertical_anchor = MSO_ANCHOR.TOP
left_content_para = left_content_frame.paragraphs[0]
left_content_para.text = ""  # Empty for user to fill
left_content_para.font.size = Pt(12)
left_content_para.font.color.rgb = RGBColor(80, 80, 80)
left_content_para.font.name = 'Segoe UI'

# Left column border
left_content_box.line.color.rgb = RGBColor(220, 220, 220)
left_content_box.line.width = Pt(1)
left_content_box.fill.background()

# Right Column: Pain Points
right_col_left = Inches(5.1)
right_col_top = Inches(4.7)
right_col_width = Inches(4.4)
right_col_height = Inches(2.5)

# Right column heading
right_heading_box = slide.shapes.add_textbox(right_col_left, right_col_top, right_col_width, Inches(0.4))
right_heading_frame = right_heading_box.text_frame
right_heading_para = right_heading_frame.paragraphs[0]
right_heading_para.text = "Pain Points"
right_heading_para.font.size = Pt(18)
right_heading_para.font.bold = True
right_heading_para.font.color.rgb = RGBColor(40, 60, 80)
right_heading_para.font.name = 'Segoe UI'

# Right column content area
right_content_box = slide.shapes.add_textbox(right_col_left, right_col_top + Inches(0.5), right_col_width, Inches(1.8))
right_content_frame = right_content_box.text_frame
right_content_frame.word_wrap = True
right_content_frame.vertical_anchor = MSO_ANCHOR.TOP
right_content_para = right_content_frame.paragraphs[0]
right_content_para.text = ""  # Empty for user to fill
right_content_para.font.size = Pt(12)
right_content_para.font.color.rgb = RGBColor(80, 80, 80)
right_content_para.font.name = 'Segoe UI'

# Right column border
right_content_box.line.color.rgb = RGBColor(220, 220, 220)
right_content_box.line.width = Pt(1)
right_content_box.fill.background()

# Save presentation
prs.save('Business_Needs_Slide.pptx')
print("✓ PowerPoint slide created successfully: Business_Needs_Slide.pptx")
print("\nSlide Features:")
print("• Word cloud (60-70% top section) - fully editable text")
print("• Two-column layout (bottom 30%)")
print("  - Left: Business Needs")
print("  - Right: Pain Points")
print("• Clean, modern design with neutral colors")
print("• All text and elements are editable in PowerPoint")
print("• Ready to download and customize!")
