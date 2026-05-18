#!/usr/bin/env python3
"""Generate Business Needs PowerPoint Slide with Word Cloud"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_business_slide():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Set background to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ===== WORD CLOUD SECTION (Top 60-70%) =====
    word_data = [
        ('PowerPoint', 50),
        ('Adobe', 45),
        ('Create', 48),
        ('Presenter', 42),
        ('Picture', 40),
        ('bulleted', 38),
        ('Content', 39),
        ('List', 37),
        ('Use', 36),
        ('Text', 35),
        ('Bullet', 34),
        ('Categories', 33),
        ('Single', 32),
        ('Learn', 31),
        ('Within', 44),
        ('Slides', 30),
        ('Add', 29),
        ('Quizzes', 28),
        ('Numbered', 27),
        ('Import', 26),
        ('Objects', 25),
        ('Custom', 24),
        ('Shape', 23),
        ('Arrange', 22),
        ('Format', 21),
        ('Question', 20),
        ('Using', 19),
        ('Insert', 18),
        ('New', 17),
        ('Read', 16),
    ]

    # Create word cloud text box
    word_cloud_left = Inches(0.5)
    word_cloud_top = Inches(0.5)
    word_cloud_width = Inches(9)
    word_cloud_height = Inches(4.2)

    word_cloud_box = slide.shapes.add_textbox(word_cloud_left, word_cloud_top, word_cloud_width, word_cloud_height)
    word_cloud_frame = word_cloud_box.text_frame
    word_cloud_frame.word_wrap = True
    word_cloud_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Normalize font sizes (20pt to 56pt range)
    max_freq = max([freq for _, freq in word_data])
    min_freq = min([freq for _, freq in word_data])

    # Add words with varying font sizes
    for i, (word, frequency) in enumerate(word_data):
        # Calculate font size based on frequency
        font_size = 20 + (frequency - min_freq) / (max_freq - min_freq) * 36
        font_size = int(font_size)
        
        if i == 0:
            p = word_cloud_frame.paragraphs[0]
        else:
            p = word_cloud_frame.add_paragraph()
        
        p.text = word
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(font_size)
        p.font.bold = (frequency > 40)
        p.font.color.rgb = RGBColor(40, 60, 80)
        p.font.name = 'Segoe UI'
        p.space_before = Pt(2)
        p.space_after = Pt(2)

    # Add border to word cloud section
    word_cloud_box.line.color.rgb = RGBColor(220, 220, 220)
    word_cloud_box.line.width = Pt(1)
    word_cloud_box.fill.background()

    # ===== TWO-COLUMN SECTION (Bottom ~30%) =====
    # Left Column: Business Needs
    left_col_left = Inches(0.5)
    left_col_top = Inches(4.7)
    left_col_width = Inches(4.4)

    # Left column heading
    left_heading_box = slide.shapes.add_textbox(left_col_left, left_col_top, left_col_width, Inches(0.4))
    left_heading_frame = left_heading_box.text_frame
    left_heading_para = left_heading_frame.paragraphs[0]
    left_heading_para.text = "Business Needs"
    left_heading_para.font.size = Pt(18)
    left_heading_para.font.bold = True
    left_heading_para.font.color.rgb = RGBColor(40, 60, 80)
    left_heading_para.font.name = 'Segoe UI'

    # Left column content area
    left_content_box = slide.shapes.add_textbox(left_col_left, left_col_top + Inches(0.5), left_col_width, Inches(1.8))
    left_content_frame = left_content_box.text_frame
    left_content_frame.word_wrap = True
    left_content_frame.vertical_anchor = MSO_ANCHOR.TOP
    left_content_para = left_content_frame.paragraphs[0]
    left_content_para.text = "" # Empty for user to fill
    left_content_para.font.size = Pt(12)
    left_content_para.font.color.rgb = RGBColor(80, 80, 80)
    left_content_para.font.name = 'Segoe UI'

    # Left column border
    left_content_box.line.color.rgb = RGBColor(220, 220, 220)
    left_content_box.line.width = Pt(1)
    left_content_box.fill.background()

    # Right Column: Pain Points
    right_col_left = Inches(5.1)
    right_col_top = Inches(4.7)
    right_col_width = Inches(4.4)

    # Right column heading
    right_heading_box = slide.shapes.add_textbox(right_col_left, right_col_top, right_col_width, Inches(0.4))
    right_heading_frame = right_heading_box.text_frame
    right_heading_para = right_heading_frame.paragraphs[0]
    right_heading_para.text = "Pain Points"
    right_heading_para.font.size = Pt(18)
    right_heading_para.font.bold = True
    right_heading_para.font.color.rgb = RGBColor(40, 60, 80)
    right_heading_para.font.name = 'Segoe UI'

    # Right column content area
    right_content_box = slide.shapes.add_textbox(right_col_left, right_col_top + Inches(0.5), right_col_width, Inches(1.8))
    right_content_frame = right_content_box.text_frame
    right_content_frame.word_wrap = True
    right_content_frame.vertical_anchor = MSO_ANCHOR.TOP
    right_content_para = right_content_frame.paragraphs[0]
    right_content_para.text = "" # Empty for user to fill
    right_content_para.font.size = Pt(12)
    right_content_para.font.color.rgb = RGBColor(80, 80, 80)
    right_content_para.font.name = 'Segoe UI'

    # Right column border
    right_content_box.line.color.rgb = RGBColor(220, 220, 220)
    right_content_box.line.width = Pt(1)
    right_content_box.fill.background()

    # Save presentation
    prs.save('Business_Needs_Slide.pptx')
    print('✓ PowerPoint slide created: Business_Needs_Slide.pptx')
    print('\nSlide Features:')
    print('• Word cloud (60-70% top) with varying font sizes')
    print('• Business Needs column (bottom left)')
    print('• Pain Points column (bottom right)')
    print('• Clean, modern design with neutral colors')
    print('• All elements fully editable in PowerPoint')

if __name__ == '__main__':
    create_business_slide()
